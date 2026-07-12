# -*- coding: utf-8 -*-
"""
基于中期汇报.md生成中期汇报PPT
风格：简洁专业（白色背景，适合投影展示）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import os

# ====== 颜色定义 ======
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
CARD_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
PRIMARY_BLUE = RGBColor(0x1E, 0x88, 0xE5)
PRIMARY_DARK = RGBColor(0x15, 0x65, 0xC0)
ACCENT_PINK = RGBColor(0xE9, 0x1E, 0x63)
ACCENT_GREEN = RGBColor(0x43, 0xA0, 0x47)
ACCENT_ORANGE = RGBColor(0xFF, 0x8F, 0x00)
ACCENT_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
TEXT_BLACK = RGBColor(0x21, 0x21, 0x21)
TEXT_GRAY = RGBColor(0x61, 0x61, 0x61)
TEXT_LIGHT = RGBColor(0x9E, 0x9E, 0x9E)
LINE_BLUE = RGBColor(0xBB, 0xDE, 0xFB)

# ====== 创建演示文稿 ======
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=22,
                color=TEXT_BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8),
                title_text, font_size=36, color=PRIMARY_DARK, bold=True)

    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.5),
                    subtitle_text, font_size=18, color=TEXT_GRAY)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4),
                                   Inches(11.7), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.fill.background()

def add_card(slide, left, top, width, height, color=CARD_LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = LINE_BLUE
    shape.line.width = Pt(1)
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=22, color=TEXT_BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(10)
    return txBox

def add_table(slide, left, top, width, height, data, header_color=PRIMARY_BLUE):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(20)
                p.font.name = '微软雅黑'
                if i == 0:
                    p.font.color.rgb = BG_WHITE
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = TEXT_BLACK
                    p.alignment = PP_ALIGN.CENTER

    for j in range(cols):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE

    for row in table.rows:
        row.height = Inches(0.7)

    return table


# ====== 第1页：封面 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
            '🐱 基于声音与行为分析的', font_size=40, color=PRIMARY_DARK, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(1.2),
            '「喵语」情感识别与需求解读系统', font_size=40, color=PRIMARY_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
            '—— 中期汇报 ——', font_size=24, color=TEXT_GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.4), Inches(11.3), Inches(0.5),
            '汇报人：李歆辰', font_size=20, color=TEXT_BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.9), Inches(11.3), Inches(0.5),
            '班级：国和中学 七（2）班', font_size=20, color=TEXT_BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
            '时间：2026年7月', font_size=20, color=TEXT_BLACK, alignment=PP_ALIGN.CENTER)

# ====== 第2页：项目背景 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '📋 项目背景', '为什么要做这个项目？')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
            '🐱 现实背景', font_size=24, color=ACCENT_PINK, bold=True)

add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.2),
                ['我家养了两只猫——可乐和草草',
                 '猫咪会通过不同叫声和肢体动作表达需求',
                 '但作为主人，经常面临「猜不透」的尴尬',
                 '希望用AI技术打破人与猫之间的沟通壁垒'],
                font_size=20)

add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.6),
            '🤖 技术契机', font_size=24, color=ACCENT_PINK, bold=True)

add_bullet_list(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.5),
                ['AI图像识别和声音分类技术已逐渐普及',
                 '百度EasyDL等平台让中学生也能使用AI',
                 '利用Trae等AI编程工具可以快速搭建应用'],
                font_size=20)

# ====== 第3页：研究思路 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🔬 研究思路与技术路线', '多模态融合')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
            '核心思路：将「理解猫语」拆解为两个子问题', font_size=22, color=PRIMARY_BLUE, bold=True)

add_card(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.0))
add_textbox(slide, Inches(1.0), Inches(2.8), Inches(5.1), Inches(0.6),
            '👀 视觉识别', font_size=24, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(1.0), Inches(3.4), Inches(5.1), Inches(1.2),
            '「它在做什么？」\n通过图片分析猫咪的姿态和动作',
            font_size=18, color=TEXT_BLACK)

add_card(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(2.0))
add_textbox(slide, Inches(7.2), Inches(2.8), Inches(5.1), Inches(0.6),
            '👂 听觉识别', font_size=24, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.2), Inches(3.4), Inches(5.1), Inches(1.2),
            '「它在怎么叫？」\n通过音频分析猫咪的叫声特征',
            font_size=18, color=TEXT_BLACK)

add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.6),
            '技术路线', font_size=22, color=ACCENT_ORANGE, bold=True)

add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.5),
            '数据采集  →  人工标注  →  AI训练  →  模型验证  →  成果集成',
            font_size=24, color=PRIMARY_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7),
            '最终结合声音和图像的分析结果，给出综合判断',
            font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ====== 第4页：情绪状态定义 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '📊 5种情绪状态定义', '通过观察建立的分类标准')

data = [
    ['编号', '状态名称', '图标', '典型表现'],
    ['①', '饥饿乞食', '🍚', '短促叫声、蹭碗、跟随主人、尾巴竖起'],
    ['②', '警惕警告', '⚠️', '飞机耳、瞳孔放大、尾巴猛甩、低沉嘶吼'],
    ['③', '心情愉悦', '😊', '呼噜声、眯眼睛、缓慢眨眼、身体放松'],
    ['④', '无聊求玩', '🎾', '叼玩具找主人、扒门、跑酷、兴奋短叫'],
    ['⑤', '标记领地', '🏠', '蹭家具、抓挠沙发、尾巴竖起颤抖'],
]
add_table(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.8), data)

add_card(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.4))
add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.2),
            '💡 设计说明：最初定义了「撒娇求关注」和「无聊求玩」两个状态，但观察发现两者表现相似容易混淆，最终用「标记领地」替换，使5种状态的区分度更高。',
            font_size=18, color=TEXT_GRAY)

# ====== 第5页：文献调研 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '📚 文献调研', '「喵语」规律的科学依据')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
            '发现「喵语」的三个规律：', font_size=22, color=PRIMARY_BLUE, bold=True)

add_card(slide, Inches(0.8), Inches(2.6), Inches(3.7), Inches(2.6))
add_textbox(slide, Inches(1.0), Inches(2.8), Inches(3.3), Inches(0.6),
            '⏱️ 叫声时长', font_size=22, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(1.0), Inches(3.5), Inches(3.3), Inches(1.6),
            '积极情绪时叫声较短\n（如吃饭、玩耍）\n\n消极情绪时叫声拖长\n（如被关进航空箱）',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(4.8), Inches(2.6), Inches(3.7), Inches(2.6))
add_textbox(slide, Inches(5.0), Inches(2.8), Inches(3.3), Inches(0.6),
            '🎵 音调高低', font_size=22, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(5.0), Inches(3.5), Inches(3.3), Inches(1.6),
            '积极情绪时音调偏高\n\n消极情绪时音调偏低\n\n这些规律为数据\n标注提供了理论基础',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(8.8), Inches(2.6), Inches(3.7), Inches(2.6))
add_textbox(slide, Inches(9.0), Inches(2.8), Inches(3.3), Inches(0.6),
            '📈 音调走势', font_size=22, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(9.0), Inches(3.5), Inches(3.3), Inches(1.6),
            '如果一声「喵」是先升高\n再降低，呈「∧」形\n\n通常是开心的表现',
            font_size=16, color=TEXT_BLACK)

# ====== 第6页：数据采集方法 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🎬 数据采集方法', '视频→关键帧的创新方法')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
            '发现问题：单张图片难以准确标注情绪', font_size=22, color=ACCENT_PINK, bold=True)

add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.8),
                ['尾巴是「缓慢摆动」还是「猛甩」？静态图片无法区分',
                 '耳朵是「慢慢变平」还是「突然贴头」？需要动态序列',
                 '同一张图片可能有多种情绪解释，容易误判'],
                font_size=20)

add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6),
            '解决方案：视频录制 → 关键帧提取', font_size=22, color=PRIMARY_BLUE, bold=True)

add_card(slide, Inches(0.8), Inches(5.3), Inches(2.6), Inches(1.6))
add_textbox(slide, Inches(0.9), Inches(5.5), Inches(2.4), Inches(1.2),
            '📹 录制视频\n（10-30秒）', font_size=18, color=ACCENT_GREEN,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(3.5), Inches(5.7), Inches(0.5), Inches(0.6),
            '→', font_size=28, color=PRIMARY_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(4.0), Inches(5.3), Inches(2.6), Inches(1.6))
add_textbox(slide, Inches(4.1), Inches(5.5), Inches(2.4), Inches(1.2),
            '👀 观看判断\n准确判断情绪', font_size=18, color=ACCENT_GREEN,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(6.7), Inches(5.7), Inches(0.5), Inches(0.6),
            '→', font_size=28, color=PRIMARY_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(7.2), Inches(5.3), Inches(2.6), Inches(1.6))
add_textbox(slide, Inches(7.3), Inches(5.5), Inches(2.4), Inches(1.2),
            '🖼️ 截取关键帧\n最典型画面', font_size=18, color=ACCENT_GREEN,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(9.9), Inches(5.7), Inches(0.5), Inches(0.6),
            '→', font_size=28, color=PRIMARY_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(10.4), Inches(5.3), Inches(2.2), Inches(1.6))
add_textbox(slide, Inches(10.5), Inches(5.5), Inches(2.0), Inches(1.2),
            '🏷️ 标注\n打上标签', font_size=18, color=ACCENT_GREEN,
            bold=True, alignment=PP_ALIGN.CENTER)

# ====== 第7页：数据采集进度 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '📈 数据采集进度', '目前采集情况')

data = [
    ['情绪状态', '可乐视频', '草草视频', '可乐音频', '草草音频', '关键帧'],
    ['饥饿乞食', '已采集', '已采集', '已采集', '已采集', '已截取'],
    ['警惕警告', '较少', '较少', '较少', '较少', '待补充'],
    ['心情愉悦', '已采集', '已采集', '已采集', '已采集', '已截取'],
    ['无聊求玩', '已采集', '较少', '已采集', '较少', '部分截取'],
    ['标记领地', '待采集', '待采集', '待采集', '待采集', '待采集'],
]
add_table(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.2), data)

add_card(slide, Inches(0.6), Inches(5.2), Inches(5.8), Inches(2.0))
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(5.4), Inches(0.6),
            '🖼️ 图片/视频帧', font_size=22, color=PRIMARY_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(6.0), Inches(5.4), Inches(1.2),
            '已截取约 120 张\n典型姿态图片', font_size=20, color=TEXT_BLACK)

add_card(slide, Inches(6.8), Inches(5.2), Inches(5.8), Inches(2.0))
add_textbox(slide, Inches(7.0), Inches(5.4), Inches(5.4), Inches(0.6),
            '🎵 音频片段', font_size=22, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(7.0), Inches(6.0), Inches(5.4), Inches(1.2),
            '已采集约 20 条\n数据量较少，待改进', font_size=20, color=TEXT_BLACK)

# ====== 第8页：开发工具 - Trae与Vibe Coding ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🛠️ 开发工具：Trae与Vibe Coding', 'AI辅助编程')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
            '什么是Vibe Coding？', font_size=24, color=PRIMARY_BLUE, bold=True)

add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
            '用自然语言描述需求，AI自动生成代码。比如我告诉Trae：\n'
            '「帮我做一个猫咪图片上传区域，支持拖拽上传，上传后显示预览图」\n'
            'Trae就能自动生成对应的HTML、CSS和JavaScript代码。',
            font_size=18, color=TEXT_BLACK)

add_card(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0))
add_textbox(slide, Inches(1.0), Inches(4.2), Inches(5.1), Inches(0.6),
            '✅ 优点', font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(4.9), Inches(5.1), Inches(2.0),
                ['专注于功能设计和逻辑思考',
                 '不用纠结于语法细节',
                 '大大提高开发效率',
                 '降低了编程门槛'],
                font_size=18)

add_card(slide, Inches(6.9), Inches(4.0), Inches(5.5), Inches(3.0))
add_textbox(slide, Inches(7.1), Inches(4.2), Inches(5.1), Inches(0.6),
            '⚠️ 注意', font_size=22, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(4.9), Inches(5.1), Inches(2.0),
                ['AI生成的代码需要自己理解',
                 '不能盲目照搬',
                 '调试和理解代码同样重要',
                 '过程中学到了很多前端知识'],
                font_size=18)

# ====== 第9页：系统架构 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🏗️ 系统架构', '平台技术架构')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
            '技术栈：原生 HTML / CSS / JavaScript（前端纯静态，无需服务器）',
            font_size=20, color=ACCENT_GREEN)

add_card(slide, Inches(0.8), Inches(2.6), Inches(6.0), Inches(4.6))
add_textbox(slide, Inches(1.0), Inches(2.8), Inches(5.6), Inches(0.6),
            '📁 项目结构', font_size=22, color=PRIMARY_BLUE, bold=True)

add_textbox(slide, Inches(1.0), Inches(3.5), Inches(5.6), Inches(3.6),
            'myfirstai/\n'
            '├── index.html\n'
            '│     # 主页面\n'
            '├── css/style.css\n'
            '│     # 赛博朋克风格样式\n'
            '├── js/\n'
            '│   ├── app.js\n'
            '│   │   # 主应用逻辑\n'
            '│   ├── model-service.js\n'
            '│   │   # 模型API调用\n'
            '│   ├── fusion.js\n'
            '│   │   # 多模态融合\n'
            '│   └── config.js\n'
            '│       # API配置\n'
            '├── assets/\n'
            '│     # 猫咪照片\n'
            '└── models/\n'
            '      # 预留模型目录',
            font_size=16, color=TEXT_BLACK, font_name='Consolas')

add_card(slide, Inches(7.2), Inches(2.6), Inches(5.3), Inches(4.6))
add_textbox(slide, Inches(7.4), Inches(2.8), Inches(4.9), Inches(0.6),
            '🔧 核心模块', font_size=22, color=PRIMARY_BLUE, bold=True)

add_bullet_list(slide, Inches(7.4), Inches(3.5), Inches(4.9), Inches(3.6),
                ['app.js - 主应用逻辑\n  界面交互、事件处理',
                 'model-service.js - 模型服务\n  EasyDL API调用、演示模式',
                 'fusion.js - 多模态融合\n  图片+音频综合判断',
                 'config.js - 配置文件\n  API密钥、平台设置'],
                font_size=16)

# ====== 第10页：多模态融合 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🔀 多模态融合', '结合图片和音频综合判断')

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.2),
            '就像我们日常生活中理解猫咪——也是同时看它的动作和听它的\n'
            '叫声来判断的。只看图片或只听声音容易误判，结合起来更准确。',
            font_size=20, color=TEXT_BLACK)

add_card(slide, Inches(0.8), Inches(3.3), Inches(3.7), Inches(2.6))
add_textbox(slide, Inches(1.0), Inches(3.5), Inches(3.3), Inches(0.6),
            '🖼️ 只有图片', font_size=20, color=ACCENT_GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(4.2), Inches(3.3), Inches(1.6),
            '只看图片结果\n\n图片权重 = 100%\n音频权重 = 0%',
            font_size=18, color=TEXT_BLACK, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(4.8), Inches(3.3), Inches(3.7), Inches(2.6))
add_textbox(slide, Inches(5.0), Inches(3.5), Inches(3.3), Inches(0.6),
            '🎵 只有音频', font_size=20, color=ACCENT_GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(5.0), Inches(4.2), Inches(3.3), Inches(1.6),
            '只听音频结果\n\n图片权重 = 0%\n音频权重 = 100%',
            font_size=18, color=TEXT_BLACK, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(8.8), Inches(3.3), Inches(3.7), Inches(2.6))
add_textbox(slide, Inches(9.0), Inches(3.5), Inches(3.3), Inches(0.6),
            '🔀 两者都有', font_size=20, color=PRIMARY_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.0), Inches(4.2), Inches(3.3), Inches(1.6),
            '综合判断\n\n图片权重 = 50%\n音频权重 = 50%',
            font_size=18, color=TEXT_BLACK, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
            '💡 系统根据用户上传内容自动调整权重，选择得分最高的状态作为最终结果',
            font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ====== 第11页：平台功能展示 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🖥️ 平台功能展示', '已实现的9项功能')

data = [
    ['功能模块', '状态', '说明'],
    ['图片上传识别', '✅', '支持拖拽/点击上传，预览，识别'],
    ['音频录制识别', '✅', '支持麦克风录音、文件上传、波形显示'],
    ['多模态融合', '✅', '图片+音频综合判断'],
    ['5种状态分类', '✅', '饥饿乞食、警惕警告、心情愉悦、无聊求玩、标记领地'],
    ['结果展示', '✅', '状态标签、置信度进度条、建议行动'],
    ['猫咪选择器', '✅', '可乐/草草/其他猫咪三选一'],
    ['历史记录', '✅', '本地存储最近10次识别记录'],
    ['演示模式', '✅', '无需真实模型即可体验完整功能'],
    ['EasyDL接口', '✅', '代码已就绪，待填入API密钥'],
]
add_table(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.2), data)

add_textbox(slide, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.4),
            '🎨 界面风格：赛博朋克科技风（深色主题 + 霓虹蓝/粉发光效果）  |  代码已推送至GitHub',
            font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ====== 第12页：演示模式 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '🎬 演示模式', '无需真实模型即可体验完整功能')

add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
add_textbox(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(0.6),
            '🎬 演示模式', font_size=24, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.9), Inches(5.1), Inches(3.2),
                ['使用模拟数据',
                 '可完整体验识别流程',
                 '适合展示项目效果',
                 '演示按钮为黄色',
                 '可乐&草草一键演示识别'],
                font_size=18)

add_card(slide, Inches(6.9), Inches(2.0), Inches(5.5), Inches(4.6))
add_textbox(slide, Inches(7.1), Inches(2.2), Inches(5.1), Inches(0.6),
            '🔍 正式模式', font_size=24, color=PRIMARY_BLUE, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(2.9), Inches(5.1), Inches(3.2),
                ['调用真实EasyDL API',
                 '需要配置API密钥',
                 '上传真实图片/音频',
                 '正式按钮为蓝色',
                 '模型训练完成后启用'],
                font_size=18)

add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
            '💡 演示功能和正式功能在界面上有明显区分，避免混淆',
            font_size=18, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# ====== 第13页：遇到的困难 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '⚠️ 遇到的困难与解决方案', '5个主要困难')

data = [
    ['困难', '解决方案'],
    ['猫咪不配合拍摄\n想要录「警告」它偏睡觉',
     '蹲守法 + 自然触发\n拿罐头→饥饿，逗猫棒→玩耍'],
    ['声音数据质量参差不齐\n有环境噪音干扰',
     '选择安静时段录制\n用Audacity剪辑3-5秒片段'],
    ['单张图片难以准确标注\n同一图片多种解释',
     '改用视频→关键帧方法\n结合上下文准确判断'],
    ['AI平台使用需要摸索\n不清楚如何设置标签',
     '看官方教学视频\n先少量数据试跑确认流程'],
    ['是否为两只猫分别训练\n纠结模型设计',
     '决定使用统一模型\n合并训练更准确更通用'],
]
add_table(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.4), data)

# ====== 第14页：下一步计划 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '📅 下一步工作计划', '7月-9月时间规划')

data = [
    ['时间', '任务', '说明'],
    ['7月中下旬', '补全缺失数据', '重点是「警惕警告」和「标记领地」，家人协同采集'],
    ['8月上旬', '完成数据标注与上传', '完成关键帧提取、音频剪辑，上传至EasyDL'],
    ['8月上旬', '启动AI模型训练', '训练声音分类器和姿态分类器'],
    ['8月中旬', '模型测试与优化', '用10%数据做测试集，准确率<80%则补充重训'],
    ['8月下旬', '模型接入平台', '配置EasyDL API，实现真正的「喵语翻译」'],
    ['9月初', '撰写结题报告', '整理研究报告，制作答辩PPT'],
]
add_table(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.2), data)

add_card(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(1.0))
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.7),
            '📌 重点：音频数据目前仅约20条，计划在采集过程中参考图片采集经验，探索如何改进音频采集流程，提高数据数量和质量。',
            font_size=18, color=ACCENT_ORANGE)

# ====== 第15页：阶段收获 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, '💡 阶段收获', '6点体会')

add_card(slide, Inches(0.6), Inches(1.8), Inches(5.9), Inches(1.7))
add_textbox(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(0.5),
            '1️⃣ 科学研究真的很辛苦', font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.45), Inches(5.5), Inches(1.0),
            '光是录猫咪叫声就花了几十个小时，\n还要反复筛选、剪辑。',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.9), Inches(1.7))
add_textbox(slide, Inches(7.0), Inches(1.95), Inches(5.5), Inches(0.5),
            '2️⃣ 理论和实践差距很大', font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.0), Inches(2.45), Inches(5.5), Inches(1.0),
            '文献规律清晰，实际采集才发现\n每只猫都有自己的「口音」和习惯。',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(0.6), Inches(3.7), Inches(5.9), Inches(1.7))
add_textbox(slide, Inches(0.8), Inches(3.85), Inches(5.5), Inches(0.5),
            '3️⃣ AI不是魔法', font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(4.35), Inches(5.5), Inches(1.0),
            '需要高质量数据才能工作，数据好坏\n直接决定模型效果。',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(6.8), Inches(3.7), Inches(5.9), Inches(1.7))
add_textbox(slide, Inches(7.0), Inches(3.85), Inches(5.5), Inches(0.5),
            '4️⃣ 「数据为王」的含义', font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.0), Inches(4.35), Inches(5.5), Inches(1.0),
            '没有干净、充足、标注准确的数据，\n再厉害的算法也跑不起来。',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(0.6), Inches(5.6), Inches(5.9), Inches(1.7))
add_textbox(slide, Inches(0.8), Inches(5.75), Inches(5.5), Inches(0.5),
            '5️⃣ 多模态融合的思路有价值', font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(6.25), Inches(5.5), Inches(1.0),
            '单一模态容易误判，结合两种模态\n能显著提升准确率。',
            font_size=16, color=TEXT_BLACK)

add_card(slide, Inches(6.8), Inches(5.6), Inches(5.9), Inches(1.7))
add_textbox(slide, Inches(7.0), Inches(5.75), Inches(5.5), Inches(0.5),
            '6️⃣ Vibe Coding让编程不一样', font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.0), Inches(6.25), Inches(5.5), Inches(1.0),
            'AI生成代码降低门槛，但仍需自己\n理解和调试，不能盲目照搬。',
            font_size=16, color=TEXT_BLACK)

# ====== 第16页：致谢 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
            '🐱 谢谢聆听！', font_size=48, color=PRIMARY_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
            '—— 期待与可乐和草草的「对话」 ——', font_size=24, color=TEXT_BLACK,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
            '感谢老师的指导！', font_size=20, color=TEXT_GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.5),
            '感谢家人的协助与支持！', font_size=20, color=TEXT_GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
            '李歆辰  |  国和中学 七（2）班  |  2026年7月',
            font_size=18, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ====== 保存文件 ======
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '中期汇报.pptx')
prs.save(output_path)
print('PPT已生成：' + output_path)
print('共 ' + str(len(prs.slides)) + ' 页')
